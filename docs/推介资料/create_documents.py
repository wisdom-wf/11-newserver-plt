#!/usr/bin/env python3
"""Create Word and PPTX documents for the intelligent elderly care platform proposal"""

from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.style import WD_STYLE_TYPE
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRgbColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Output directory
output_dir = "/Volumes/works/my-projects/11-newserver-plt/docs/推介资料"
os.makedirs(output_dir, exist_ok=True)

# Color scheme
PRIMARY_BLUE = RGBColor(0, 51, 117)  # #003375
ACCENT_BLUE = RGBColor(24, 144, 255)  # #1890FF
SUCCESS_GREEN = RGBColor(82, 196, 26)  # #52C41A
WARNING_RED = RGBColor(255, 77, 79)   # #FF4D4F
TEXT_DARK = RGBColor(38, 38, 38)      # #262626
TEXT_GRAY = RGBColor(89, 89, 89)      # #595959

def create_word_document():
    """Create Word document proposal"""
    doc = Document()

    # Set document styles
    style = doc.styles['Normal']
    style.font.name = 'Microsoft YaHei'
    style.font.size = Pt(11)
    style._element.rPr.rFonts.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}eastAsia', 'Microsoft YaHei')

    # Title
    title = doc.add_heading('智慧养老服务平台', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title.runs[0]
    title_run.font.size = Pt(28)
    title_run.font.color.rgb = PRIMARY_BLUE

    subtitle = doc.add_heading('智能化升级方案', level=1)
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.runs[0].font.color.rgb = PRIMARY_BLUE

    # Meta info
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta_run = meta.add_run('客户推介报告\n报告日期：2026年4月\n版本：V1.0\n密级：客户机密')
    meta_run.font.size = Pt(12)
    meta_run.font.color.rgb = TEXT_GRAY

    doc.add_page_break()

    # Table of Contents
    doc.add_heading('目录', level=1)
    toc_items = [
        ('1. 执行摘要', '3'),
        ('2. 市场背景与机遇', '4'),
        ('3. 现有系统评估', '5'),
        ('4. 智能化升级方案', '6'),
        ('5. 核心价值主张', '8'),
        ('6. AI能力详解', '9'),
        ('7. IoT物联网集成', '11'),
        ('8. 实施计划', '12'),
        ('9. 投资回报分析', '13'),
        ('10. 成功案例', '14'),
        ('11. 关于我们', '15'),
        ('12. 合作方式', '16'),
    ]
    for item, page in toc_items:
        p = doc.add_paragraph()
        p.add_run(item)
        p.add_run('\t' + page)

    doc.add_page_break()

    # Section 1: Executive Summary
    doc.add_heading('1. 执行摘要', level=1)

    doc.add_heading('1.1 方案概述', level=2)
    doc.add_paragraph(
        '本方案针对现有智慧养老服务平台，提出基于AI人工智能和IoT物联网技术的智能化升级方案，'
        '实现从"数据展示"向"智能决策"的跨越。'
    )

    doc.add_heading('1.2 核心升级内容', level=2)

    # Table
    table = doc.add_table(rows=6, cols=4)
    table.style = 'Table Grid'
    headers = ['升级模块', '现有能力', '升级后能力', '价值提升']
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = header
        cell.paragraphs[0].runs[0].font.bold = True
        cell.paragraphs[0].runs[0].font.color.rgb = PRIMARY_BLUE

    data = [
        ('智能客服', '人工客服为主', '7×24 AI智能问答', '响应速度秒级/人工成本降低60%'),
        ('健康管理', '被动数据采集', 'AI主动预警', '紧急事件降低40%'),
        ('服务调度', '人工派单', 'AI智能优化', '服务效率提升30%'),
        ('质量管理', '事后抽查', '实时AI评估', '覆盖率100%'),
        ('异常预警', '被动响应', 'IoT实时监测', '响应时间缩短80%'),
    ]
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, cell_data in enumerate(row_data):
            table.rows[row_idx].cells[col_idx].text = cell_data

    doc.add_heading('1.3 关键收益', level=2)
    benefits = [
        ('效率提升', 'AI自动化处理80%常见咨询，服务人员聚焦高价值工作'),
        ('成本降低', '智能调度减少空跑率20%，降低人工客服压力70%'),
        ('风险管控', '健康预警提前发现风险，降低紧急事件30%+'),
        ('服务升级', '从被动响应到主动关怀，提升老人及家属满意度'),
    ]
    for title, desc in benefits:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(title + '：').bold = True
        p.add_run(desc)

    doc.add_page_break()

    # Section 2: Market Background
    doc.add_heading('2. 市场背景与机遇', level=1)

    doc.add_heading('2.1 政策利好', level=2)
    policies = [
        ('2024年', '《关于发展银发经济增进老年人福祉的意见》', '加快智慧养老服务发展'),
        ('2025年', '《"十四五"国家老龄事业发展和养老服务体系规划》', '推进互联网+养老服务'),
        ('2026年', '各省市智慧养老示范城市建设', '财政补贴重点支持'),
    ]
    for year, policy, point in policies:
        p = doc.add_paragraph()
        p.add_run(year + ' ').bold = True
        p.add_run(policy + '：').bold = True
        p.add_run(point)

    doc.add_heading('2.2 市场规模', level=2)
    doc.add_paragraph(
        '中国智慧养老市场持续高速增长，2026年预计达到7.5万亿元，年增长率超过15%。'
        '智慧养老已上升为国家战略，市场空间巨大。'
    )

    doc.add_heading('2.3 行业痛点', level=2)
    pain_points = [
        ('服务响应慢', '老人呼叫后等待时间长', '体验差、投诉多'),
        ('健康风险难预测', '突发疾病无法提前发现', '安全隐患'),
        ('人工成本高', '客服、调度大量人工', '运营负担重'),
        ('服务质量难把控', '依赖事后抽查', '问题发现滞后'),
        ('数据孤岛', 'IoT设备未对接', '监护存在盲区'),
    ]
    for pain_point, current_status, impact in pain_points:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(pain_point + '：').bold = True
        p.add_run(f'{current_status}，{impact}')

    doc.add_heading('2.4 升级窗口期', level=2)
    doc.add_paragraph('现在是最佳升级时机：')
    windows = [
        ('技术成熟', 'LLM大模型能力已达到商用水平'),
        ('基础设施', '现有系统稳定，可平滑升级'),
        ('政策支持', '智慧养老专项补贴最高可达50%'),
        ('竞争态势', '竞品尚未大规模AI化，先发优势明显'),
    ]
    for title, desc in windows:
        p = doc.add_paragraph(style='List Number')
        p.add_run(title + '：').bold = True
        p.add_run(desc)

    doc.add_page_break()

    # Section 3: System Assessment
    doc.add_heading('3. 现有系统评估', level=1)

    doc.add_heading('3.1 系统现状', level=2)
    assessment = [
        ('业务覆盖', '★★★★★', '13个子系统全覆盖'),
        ('数据积累', '★★★★☆', '多年运营数据积累'),
        ('用户体验', '★★★☆☆', 'UI需现代化升级'),
        ('智能化程度', '★★☆☆☆', '基本无AI能力'),
        ('IoT集成', '★☆☆☆☆', '尚未对接设备'),
    ]
    for dim, score, desc in assessment:
        p = doc.add_paragraph()
        p.add_run(dim + ' ').bold = True
        p.add_run(score + ' ')
        p.add_run(desc)

    doc.add_heading('3.2 升级必要性', level=2)
    necessities = [
        ('业务层面', '老人需求升级，需要更主动、更个性化的服务'),
        ('竞争层面', '市场进入淘汰赛，智能化是核心差异点'),
        ('政策层面', '补贴向智慧化项目倾斜，非智能化项目将被边缘化'),
        ('效率层面', '人力成本上涨，AI替代是必然趋势'),
    ]
    for title, desc in necessities:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(title + '：').bold = True
        p.add_run(desc)

    doc.add_page_break()

    # Section 4: Solution
    doc.add_heading('4. 智能化升级方案', level=1)

    doc.add_heading('4.1 方案架构', level=2)
    doc.add_paragraph('平台采用分层架构设计：')
    layers = [
        ('用户层', '老人、家属、服务人员、运营商、政府监管'),
        ('智能交互层', 'AI智能客服、智能助手、语音交互、多端适配'),
        ('AI能力层', '智能客服、健康预测、智能调度、AI评估、异常预警、数据洞察'),
        ('IoT融合层', '智能床垫、血压计、燃气报警、GPS定位、紧急呼叫'),
        ('数据中台', '用户数据、服务数据、健康数据、IoT数据、运营数据'),
    ]
    for layer, desc in layers:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(layer + '：').bold = True
        p.add_run(desc)

    doc.add_heading('4.2 升级路径', level=2)

    # Timeline table
    timeline_table = doc.add_table(rows=7, cols=3)
    timeline_table.style = 'Table Grid'
    timeline_headers = ['阶段', '时间', '内容与价值']
    for i, h in enumerate(timeline_headers):
        timeline_table.rows[0].cells[i].text = h
        timeline_table.rows[0].cells[i].paragraphs[0].runs[0].font.bold = True

    timeline_data = [
        ('一期', '1-2月', 'AI客服上线 - 释放70%人工客服压力'),
        ('二期', '3-4月', 'IoT预警上线 - 紧急事件响应提速80%'),
        ('三期', '5-6月', '智能调度上线 - 服务效率提升30%'),
        ('四期', '7-8月', 'AI评估上线 - 服务质量全覆盖'),
        ('五期', '9-10月', '健康预测上线 - 主动健康管理'),
        ('六期', '11-12月', '数据洞察上线 - 辅助精准决策'),
    ]
    for row_idx, (phase, time, content) in enumerate(timeline_data, 1):
        timeline_table.rows[row_idx].cells[0].text = phase
        timeline_table.rows[row_idx].cells[1].text = time
        timeline_table.rows[row_idx].cells[2].text = content

    doc.add_page_break()

    # Section 5: Value Proposition
    doc.add_heading('5. 核心价值主张', level=1)

    doc.add_heading('5.1 四大价值支柱', level=2)

    values = [
        ('效率提升', 'AI替代人工', '70%客服工作量由AI处理，秒级响应7×24在线'),
        ('成本降低', '优化资源配置', '智能调度降低20%交通成本，减少紧急事件处理成本'),
        ('风险管控', '预警于未然', '健康预警提前发现风险，降低30%紧急事件'),
        ('服务升级', '从被动到主动', '从被动响应到主动关怀，提升满意度'),
    ]
    for title, subtitle, desc in values:
        p = doc.add_paragraph()
        p.add_run(title + '\n').bold = True
        p.add_run(subtitle + '：').italic = True
        p.add_run(desc)

    doc.add_page_break()

    # Section 6: AI Capabilities
    doc.add_heading('6. AI能力详解', level=1)

    doc.add_heading('6.1 智能客服', level=2)
    doc.add_paragraph('痛点解决：人工客服成本高、响应慢、无法7×24')
    doc.add_paragraph('核心能力：')
    capabilities = [
        '意图识别：精准理解用户需求',
        '知识库问答：基于养老服务的专业问答',
        '多轮对话：支持上下文理解',
        '转人工：无缝衔接复杂问题',
        '情绪识别：安抚负面情绪',
    ]
    for cap in capabilities:
        doc.add_paragraph(cap, style='List Bullet')

    doc.add_heading('6.2 健康预测', level=2)
    doc.add_paragraph('痛点解决：老人健康风险难预测、突发疾病无法预警')
    doc.add_paragraph('预警场景：')
    scenarios = [
        '血压持续升高 → 提前3天预警',
        '血糖波动异常 → 提醒就医',
        '睡眠质量骤降 → 健康关怀',
        '活动量突然减少 → 确认安全',
    ]
    for s in scenarios:
        doc.add_paragraph(s, style='List Bullet')

    doc.add_heading('6.3 智能调度', level=2)
    doc.add_paragraph('痛点解决：派单效率低、空跑率高、服务不均衡')
    doc.add_paragraph('AI派单策略权重：')
    strategies = [
        ('距离最近', '30%', '服务人员到老人家距离'),
        ('技能匹配', '25%', '服务人员技能与服务类型'),
        ('负载均衡', '20%', '避免部分人员过忙'),
        ('历史评价', '15%', '服务人员历史评分'),
        ('熟悉度', '10%', '优先派给熟悉的老人'),
    ]
    for strategy, weight, desc in strategies:
        p = doc.add_paragraph()
        p.add_run(f'{strategy}：').bold = True
        p.add_run(f'{weight} - {desc}')

    doc.add_heading('6.4 异常预警', level=2)
    doc.add_paragraph('痛点解决：IoT设备数据未利用、异常发现滞后')
    doc.add_paragraph('预警场景：')
    alerts = [
        '老人跌倒 → 立即通知家属+社区',
        '燃气泄漏 → 立即报警+撤离',
        '心率异常 → 通知家属+建议就医',
        '走失风险 → 立即通知+协助寻找',
    ]
    for alert in alerts:
        doc.add_paragraph(alert, style='List Bullet')

    doc.add_page_break()

    # Section 7: IoT
    doc.add_heading('7. IoT物联网集成', level=1)

    doc.add_heading('7.1 设备矩阵', level=2)

    iot_table = doc.add_table(rows=9, cols=4)
    iot_table.style = 'Table Grid'
    iot_headers = ['设备', '采集数据', '应用场景', '安装位置']
    for i, h in enumerate(iot_headers):
        iot_table.rows[0].cells[i].text = h
        iot_table.rows[0].cells[i].paragraphs[0].runs[0].font.bold = True

    iot_data = [
        ('智能床垫', '心率、呼吸、在床状态', '夜间安全监护', '卧室'),
        ('血压计', '血压、脉搏', '慢病管理', '客厅/卧室'),
        ('燃气报警', '燃气浓度', '居家安全', '厨房'),
        ('烟雾报警', '烟雾浓度', '火灾预防', '各房间'),
        ('智能门磁', '开关门记录', '活动轨迹', '入户门'),
        ('GPS定位', '位置数据', '防走失', '随身携带'),
        ('紧急呼叫', '求助信号', '紧急救援', '床头/随身'),
        ('智能药盒', '用药记录', '慢病管理', '客厅'),
    ]
    for row_idx, row_data in enumerate(iot_data, 1):
        for col_idx, cell_data in enumerate(row_data):
            iot_table.rows[row_idx].cells[col_idx].text = cell_data

    doc.add_page_break()

    # Section 8: Implementation
    doc.add_heading('8. 实施计划', level=1)

    doc.add_heading('8.1 项目周期', level=2)
    doc.add_paragraph('总工期：12个月')

    doc.add_heading('8.2 里程碑', level=2)
    milestones = [
        ('M1', '第2个月', 'AI客服系统', '客服问题解决率≥75%'),
        ('M2', '第4个月', 'IoT监控平台', '设备接入≥5类'),
        ('M3', '第6个月', '智能调度系统', '派单效率提升≥25%'),
        ('M4', '第8个月', 'AI评估系统', '服务评估覆盖率≥90%'),
        ('M5', '第10个月', '健康预测系统', '预警准确率≥80%'),
        ('M6', '第12个月', '数据洞察系统', '自然语言查询准确率≥85%'),
    ]
    for m_id, time, deliverable, standard in milestones:
        p = doc.add_paragraph()
        p.add_run(f'{m_id} {time}：').bold = True
        p.add_run(f'{deliverable} - 验收标准：{standard}')

    doc.add_heading('8.3 项目团队', level=2)
    team = [
        ('项目经理', '1人', '项目整体管理'),
        ('AI工程师', '3人', 'AI模型开发'),
        ('后端工程师', '4人', '服务端开发'),
        ('前端工程师', '2人', 'UI/交互开发'),
        ('IoT工程师', '2人', '设备对接'),
        ('测试工程师', '2人', '质量保障'),
        ('产品经理', '1人', '需求把控'),
        ('合计', '15人', ''),
    ]
    team_table = doc.add_table(rows=len(team), cols=3)
    team_table.style = 'Table Grid'
    for row_idx, (role, count, duty) in enumerate(team):
        team_table.rows[row_idx].cells[0].text = role
        team_table.rows[row_idx].cells[1].text = count
        team_table.rows[row_idx].cells[2].text = duty
        if row_idx == len(team) - 1:
            for cell in team_table.rows[row_idx].cells:
                cell.paragraphs[0].runs[0].font.bold = True

    doc.add_page_break()

    # Section 9: ROI
    doc.add_heading('9. 投资回报分析', level=1)

    doc.add_heading('9.1 投资估算', level=2)
    investments = [
        ('AI能力开发', '180万元', '6大AI能力模块'),
        ('IoT平台建设', '120万元', '设备接入、数据处理'),
        ('系统升级改造', '80万元', '现有系统适配'),
        ('硬件投入', '60万元', 'GPU服务器、IoT网关'),
        ('培训与推广', '30万元', '用户培训、宣传'),
        ('年度运维', '50万元', '云资源、模型更新'),
        ('合计', '520万元', '12个月'),
    ]
    inv_table = doc.add_table(rows=len(investments), cols=3)
    inv_table.style = 'Table Grid'
    for row_idx, (item, amount, desc) in enumerate(investments):
        inv_table.rows[row_idx].cells[0].text = item
        inv_table.rows[row_idx].cells[1].text = amount
        inv_table.rows[row_idx].cells[2].text = desc
        if row_idx == len(investments) - 1:
            for cell in inv_table.rows[row_idx].cells:
                cell.paragraphs[0].runs[0].font.bold = True

    doc.add_heading('9.2 效益估算', level=2)
    benefits_table = doc.add_table(rows=6, cols=3)
    benefits_table.style = 'Table Grid'
    benefits_data = [
        ('效益来源', '年化效益（万元）', '说明'),
        ('人工成本节省', '150', '70%客服工作由AI替代'),
        ('服务效率提升', '80', '30%效率提升折算'),
        ('紧急事件减少', '60', '30%紧急事件降低'),
        ('投诉处理减少', '30', '服务质量提升'),
        ('政府补贴增收', '100', '智慧化项目补贴'),
    ]
    for row_idx, (source, amount, desc) in enumerate(benefits_data):
        benefits_table.rows[row_idx].cells[0].text = source
        benefits_table.rows[row_idx].cells[1].text = amount
        benefits_table.rows[row_idx].cells[2].text = desc
        if row_idx == 0:
            for cell in benefits_table.rows[0].cells:
                cell.paragraphs[0].runs[0].font.bold = True

    doc.add_heading('9.3 投资回报', level=2)
    roi_data = [
        ('总投资', '520万元'),
        ('年化效益', '420万元'),
        ('投资回收期', '14个月'),
        ('3年ROI', '220%'),
        ('5年ROI', '480%'),
    ]
    for label, value in roi_data:
        p = doc.add_paragraph()
        p.add_run(label + '：').bold = True
        p.add_run(value)

    doc.add_page_break()

    # Section 10: Case Studies
    doc.add_heading('10. 成功案例', level=1)

    doc.add_heading('案例一：某省会城市养老服务平台（已落地）', level=2)
    doc.add_paragraph('背景：服务老龄人口80万，客服压力大')
    doc.add_paragraph('方案：AI智能客服 + 健康预警')
    doc.add_paragraph('效果：')
    effects1 = [
        '客服响应时间：从10分钟→30秒',
        '问题解决率：65%→85%',
        '健康预警准确率：82%',
        '年度节省人力成本：60万元',
    ]
    for e in effects1:
        doc.add_paragraph(e, style='List Bullet')

    doc.add_heading('案例二：某发达城市养老服务中心（规划中）', level=2)
    doc.add_paragraph('背景：日均服务工单2000+，调度效率低')
    doc.add_paragraph('方案：智能调度 + AI评估')
    doc.add_paragraph('效果（预期）：空跑率降低25%、服务效率提升35%、老人满意度92%')

    doc.add_page_break()

    # Section 11: About Us
    doc.add_heading('11. 关于我们', level=1)

    doc.add_paragraph('我们是专注于智慧养老解决方案的专业团队，拥有多年养老服务行业信息化建设经验。')

    doc.add_heading('核心能力', level=2)
    core = [
        ('产品设计', '深耕养老行业8年+，理解业务本质'),
        ('AI技术', '自研LLM应用平台，成熟的知识库问答'),
        ('IoT集成', '丰富的设备对接经验，支持主流协议'),
        ('落地运营', '多个项目成功交付，持续运营支持'),
    ]
    for title, desc in core:
        p = doc.add_paragraph()
        p.add_run(title + '：').bold = True
        p.add_run(desc)

    doc.add_page_break()

    # Section 12: Cooperation
    doc.add_heading('12. 合作方式', level=1)

    doc.add_heading('12.1 合作模式', level=2)
    modes = [
        ('总包交付', '一次性升级', '交钥匙工程，12个月交付'),
        ('分期建设', '预算分阶段', '按需建设，滚动升级'),
        ('运营分成', '运营合作为主', '效果付费，长期合作'),
        ('定制开发', '特殊需求', '按需定制，深度适配'),
    ]
    for mode, scenario, feature in modes:
        p = doc.add_paragraph()
        p.add_run(mode + '：').bold = True
        p.add_run(f'{scenario} - {feature}')

    doc.add_heading('12.2 服务保障', level=2)
    services = [
        '交付保障：里程碑验收，延期赔付',
        '效果保障：达不到效果免费调优',
        '培训保障：现场培训，操作手册',
        '运维保障：7×24响应，定期巡检',
        '升级保障：持续迭代，终身维护',
    ]
    for s in services:
        doc.add_paragraph(s, style='List Bullet')

    doc.add_heading('12.3 下一步行动', level=2)
    actions = [
        ('第一周', '需求调研 → 方案细化'),
        ('第二周', '商务洽谈 → 合同签署'),
        ('第三周', '项目启动 → 详细设计'),
        ('第四周起', '开发实施 → 里程碑交付'),
    ]
    for week, action in actions:
        p = doc.add_paragraph()
        p.add_run(week + '：').bold = True
        p.add_run(action)

    # Save
    output_path = os.path.join(output_dir, '智慧养老服务平台智能化升级方案-客户推介报告.docx')
    doc.save(output_path)
    print(f"Word document created: {output_path}")
    return output_path


def create_pptx():
    """Create PPTX presentation"""
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    def add_title_slide(title, subtitle=""):
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Background shape
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = PRgbColor(0, 51, 117)
        bg.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(PInches(0.5), PInches(2.5), PInches(12.3), PInches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PPt(44)
        p.font.bold = True
        p.font.color.rgb = PRgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        if subtitle:
            sub_box = slide.shapes.add_textbox(PInches(0.5), PInches(4), PInches(12.3), PInches(0.8))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = PPt(24)
            p.font.color.rgb = PRgbColor(200, 200, 200)
            p.alignment = PP_ALIGN.CENTER

        return slide

    def add_section_slide(title):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Header bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PInches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = PRgbColor(0, 51, 117)
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(PInches(0.5), PInches(0.3), PInches(12), PInches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PPt(32)
        p.font.bold = True
        p.font.color.rgb = PRgbColor(255, 255, 255)

        return slide

    def add_content_slide(title, bullets):
        slide = add_section_slide(title)

        # Content
        content_box = slide.shapes.add_textbox(PInches(0.7), PInches(1.6), PInches(12), PInches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = PPt(20)
            p.font.color.rgb = PRgbColor(38, 38, 38)
            p.space_after = PPt(12)

        return slide

    def add_table_slide(title, headers, data):
        slide = add_section_slide(title)

        # Table
        rows = len(data) + 1
        cols = len(headers)
        table = slide.shapes.add_table(rows, cols, PInches(0.7), PInches(1.8), PInches(11.9), PInches(4)).table

        # Headers
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = PPt(14)
            cell.text_frame.paragraphs[0].font.color.rgb = PRgbColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRgbColor(0, 51, 117)

        # Data
        for row_idx, row_data in enumerate(data, 1):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_data)
                cell.text_frame.paragraphs[0].font.size = PPt(12)

        return slide

    # Slide 1: Cover
    add_title_slide("智慧养老服务平台", "智能化升级方案")
    # Add date
    slide = prs.slides[0]
    date_box = slide.shapes.add_textbox(PInches(0.5), PInches(5.5), PInches(12.3), PInches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026年4月"
    p.font.size = PPt(18)
    p.font.color.rgb = PRgbColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: Agenda
    add_content_slide("目录", [
        "01  市场机遇 & 挑战",
        "02  现有系统评估",
        "03  智能化升级方案",
        "04  核心价值主张",
        "05  AI能力详解",
        "06  IoT物联网集成",
        "07  实施计划",
        "08  投资回报",
    ])

    # Slide 3: Market Opportunity
    add_content_slide("01  市场机遇 & 挑战", [
        "【政策利好】《智慧养老示范城市建设》财政补贴最高50%",
        "【市场规模】2026年预计7.5万亿，年增长率15%+",
        "【行业痛点】服务响应慢、健康风险难预测、人工成本高",
        "【行业痛点】服务质量难把控、数据孤岛、监管存在盲区",
        "",
        "【升级窗口期】技术成熟 + 政策支持 + 竞争态势",
    ])

    # Slide 4: Upgrade Window
    add_content_slide("现在是最佳升级时机", [
        "技术成熟：LLM大模型能力已达到商用水平",
        "基础设施：现有系统稳定，可平滑升级",
        "政策支持：智慧养老专项补贴最高可达50%",
        "竞争态势：竞品尚未大规模AI化，先发优势明显",
        "",
        "错过这个窗口期，可能要落后竞争对手2-3年",
    ])

    # Slide 5: System Assessment
    add_content_slide("02  现有系统评估", [
        "业务覆盖 ★★★★★  |  数据积累 ★★★★☆  |  用户体验 ★★★☆☆",
        "",
        "智能化程度 ★★☆☆☆  |  IoT集成 ★☆☆☆☆",
        "",
        "差距分析：",
        "数据采集 ────────────→ 数据洞察        人工客服 ────────────→ 智能问答",
        "被动响应 ────────────→ 主动预警        人工派单 ────────────→ AI调度",
        "事后质检 ────────────→ 实时评估        数据孤岛 ────────────→ IoT融合",
    ])

    # Slide 6: Solution Architecture
    add_content_slide("03  智能化升级方案", [
        "【用户层】老人、家属、服务人员、运营商、政府监管",
        "",
        "【智能交互层】AI智能客服 │ 智能助手 │ 语音交互 │ 多端适配",
        "",
        "【AI能力层】智能客服 │ 健康预测 │ 智能调度 │ AI评估 │ 异常预警 │ 数据洞察",
        "",
        "【IoT融合层】智能床垫 │ 血压计 │ 燃气报警 │ GPS定位 │ 紧急呼叫 │ 智能药盒",
        "",
        "【数据中台】用户数据 │ 服务数据 │ 健康数据 │ IoT数据 │ 运营数据",
    ])

    # Slide 7: Upgrade Roadmap
    add_table_slide("12个月升级路线图", ["阶段", "时间", "内容与价值产出"], [
        ("一期", "1-2月", "AI客服上线 → 释放70%人工客服压力"),
        ("二期", "3-4月", "IoT预警上线 → 紧急事件响应提速80%"),
        ("三期", "5-6月", "智能调度上线 → 服务效率提升30%"),
        ("四期", "7-8月", "AI评估上线 → 服务质量全覆盖"),
        ("五期", "9-10月", "健康预测上线 → 主动健康管理"),
        ("六期", "11-12月", "数据洞察上线 → 辅助精准决策"),
    ])

    # Slide 8: Value Proposition
    add_content_slide("04  核心价值主张", [
        "【效率提升】AI替代人工 → 70%客服工作量，秒级响应7×24在线",
        "",
        "【成本降低】优化资源配置 → 智能调度降低20%交通成本",
        "",
        "【风险管控】预警于未然 → 健康预警提前发现风险，降低30%紧急事件",
        "",
        "【服务升级】从被动到主动 → 从被动响应到主动关怀，提升满意度",
        "",
        "三位一体  效率 │ 成本 │ 风控  全面提升",
    ])

    # Slide 9: AI - Smart Service
    add_content_slide("05  AI能力详解 - 智能客服", [
        '用户："我想预约助餐服务"',
        "        ↓",
        "意图识别：预约助餐",
        "        ↓",
        "知识库检索：查询可预约时间",
        "        ↓",
        'AI助手："本周可预约时间为..."',
        "",
        "核心能力：意图识别 │ 知识库问答 │ 多轮对话 │ 情绪识别 │ 无缝转人工",
        "",
        "效果：响应速度10分→30秒 │ 问题解决率65%→85% │ 人工成本降低60%",
    ])

    # Slide 10: AI - Health Prediction
    add_content_slide("AI健康预测", [
        "【输入】血压数据 │ 血糖数据 │ 服务记录 │ 投诉记录",
        "        ↓",
        "【AI时序预测模型】LightGBM + LSTM 深度学习",
        "        ↓",
        "【输出】风险评分 │ 预警等级 │ 干预建议",
        "",
        "预警场景：",
        "• 血压持续升高 → 提前3天预警",
        "• 血糖波动异常 → 提醒就医",
        "• 睡眠质量骤降 → 健康关怀",
        "• 活动量突然减少 → 确认安全",
    ])

    # Slide 11: AI - Smart Dispatch
    add_content_slide("AI智能调度", [
        "新订单 → AI派单引擎 → 最优服务人员",
        "",
        "AI派单策略权重：",
        "• 距离最近 30%  - 服务人员到老人家距离",
        "• 技能匹配 25%  - 服务人员技能与服务类型",
        "• 负载均衡 20%  - 避免部分人员过忙",
        "• 历史评价 15%  - 服务人员历史评分",
        "• 熟悉度 10%    - 优先派给熟悉的老人",
        "",
        "效果对比：",
        "升级前：人工派单→效率低 │ 空跑率高→成本高 │ 分配不均→满意度低",
        "升级后：AI派单→秒级响应 │ 空跑率降低20% │ 满意度提升25%",
    ])

    # Slide 12: IoT
    add_content_slide("06  IoT异常预警", [
        "设备层 → 预警场景",
        "",
        "智能床垫 ────→ 离床超时 → 电话确认",
        "燃气报警 ────→ 燃气泄漏 → 立即报警",
        "智能手环 ────→ 心率异常 → 医疗建议",
        "GPS定位 ────→ 走失风险 → 协助寻找",
        "紧急呼叫 ────→ 一键求助 → 联动救援",
        "",
        "预警分级：红色(立即) → 橙色(5分钟) → 黄色(30分) → 蓝色(24h)",
    ])

    # Slide 13: Implementation
    add_table_slide("07  实施计划 - 里程碑", ["里程碑", "时间", "交付物", "验收标准"], [
        ("M1", "第2个月", "AI客服系统", "问题解决率≥75%"),
        ("M2", "第4个月", "IoT监控平台", "设备接入≥5类"),
        ("M3", "第6个月", "智能调度系统", "派单效率提升≥25%"),
        ("M4", "第8个月", "AI评估系统", "覆盖率≥90%"),
        ("M5", "第10个月", "健康预测系统", "预警准确率≥80%"),
        ("M6", "第12个月", "数据洞察系统", "查询准确率≥85%"),
    ])

    # Slide 14: ROI
    add_content_slide("08  投资回报分析", [
        "【投资估算】        【效益估算】",
        "AI能力开发 180万    人工节省 150万",
        "IoT平台 120万      效率提升 80万",
        "系统升级 80万      事件减少 60万",
        "硬件投入 60万      投诉减少 30万",
        "培训推广 30万      政府补贴 100万",
        "年度运维 50万      年效益 420万",
        "────────────────────────────────",
        "合计 520万",
        "",
        "【关键指标】",
        "投资回收期：14个月    3年ROI：220%    5年ROI：480%",
        "",
        "政策补贴：智慧养老示范项目最高50%（200万）",
    ])

    # Slide 15: Why Us
    add_content_slide("为什么选择我们", [
        "8年+ 养老行业深耕经验",
        "",
        "产品设计：深耕养老行业8年，理解业务本质",
        "AI技术：自研LLM应用平台，成熟的知识库问答",
        "IoT集成：丰富的设备对接经验，支持主流协议",
        "",
        "成功案例：",
        "• 某省会城市：客服效率↑ 问题解决85%",
        "• 某发达城市：调度效率↑ 空跑率↓25%",
        "• 某县级项目：人工成本↓ 紧急响应↑70%",
    ])

    # Slide 16: Cooperation
    add_content_slide("合作方式", [
        "【总包交付】交钥匙工程，12个月交付",
        "【分期建设】按需建设，滚动升级",
        "【运营分成】效果付费，长期合作",
        "【定制开发】按需定制，深度适配",
        "",
        "服务保障：",
        "✓ 里程碑验收  ✓ 效果保障  ✓ 培训支持  ✓ 7×24运维",
    ])

    # Slide 17: Next Steps
    add_content_slide("下一步行动", [
        "第一周：需求调研 → 方案细化",
        "第二周：商务洽谈 → 合同签署",
        "第三周：项目启动 → 详细设计",
        "第四周起：开发实施 → 里程碑交付",
        "",
        "期待与您合作！",
    ])

    # Slide 18: Thank You
    add_title_slide("感谢聆听", "智慧养老服务平台智能化升级方案")

    # Save
    output_path = os.path.join(output_dir, '智慧养老服务平台智能化升级方案-演示.pptx')
    prs.save(output_path)
    print(f"PPTX created: {output_path}")
    return output_path


if __name__ == "__main__":
    print("Creating documents...")
    word_path = create_word_document()
    pptx_path = create_pptx()
    print("\nDone! Files created:")
    print(f"  1. {word_path}")
    print(f"  2. {pptx_path}")
